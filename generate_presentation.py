#!/usr/bin/env python3
"""
Monthly Executive Presentation Generator

Reads the generated project reports from outputs/, feeds them to Gemini to synthesize
portfolio trends and risks, and programmatically generates a professional 5-7 slide
PowerPoint presentation using python-pptx.
"""

import json
import os
import sys
from pathlib import Path
from google import genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Colors for professional design
COLOR_PRIMARY = RGBColor(16, 44, 87)    # Navy Blue
COLOR_SECONDARY = RGBColor(53, 162, 235) # Light Blue
COLOR_DARK = RGBColor(33, 37, 41)       # Charcoal
COLOR_MUTED = RGBColor(108, 117, 125)   # Muted Gray
COLOR_GREEN = RGBColor(40, 167, 69)     # Green
COLOR_AMBER = RGBColor(255, 193, 7)     # Amber
COLOR_RED = RGBColor(220, 53, 69)       # Red


def _get_client() -> genai.Client:
    api_key = os.environ.get("GEMINI_API_KEY")
    if not api_key:
        raise ValueError("GEMINI_API_KEY environment variable is not set.")
    return genai.Client(api_key=api_key)


def get_synthesis(reports: list) -> dict:
    client = _get_client()

    prompt = f"""You are a Principal Delivery Analyst preparing a monthly executive portfolio review presentation for a VP to present to a client.
Analyze the following project health reports and synthesize them into a structured 6-slide presentation. Do NOT just summarize each project; synthesize trends, highlight cross-project dependencies, identify systemic risks, and offer strategic recommendations.

Project Reports:
{json.dumps(reports, indent=2)}

You must return your output strictly in JSON format. Do not include markdown code block formatting (e.g. do not wrap the JSON in ```json...```). Return ONLY the raw JSON string.

The JSON structure must match this schema:
{{
  "slides": [
    {{
      "title": "Slide Title",
      "subtitle": "Optional slide subtitle",
      "layout_type": "title_slide | bullet_slide | two_column_slide",
      "content": [
         "Point 1",
         "Point 2",
         "Point 3",
         "Point 4"
      ],
      "secondary_content": [] // Used only if layout_type is two_column_slide
    }}
  ]
}}

Generate exactly 6 slides:
1. Title slide: Executive Portfolio Review (VP presenting to Client)
2. Portfolio Overview & Health Dashboard: Summary of RAG distributions (Green, Amber, Red) and overall confidence
3. Core Trends & Insights: Common patterns (e.g., strong budget discipline vs schedule slippage)
4. Key Risks & Delivery Bottlenecks: Systemic issues like vendor delays, client resource constraints
5. Resource & Financial Optimization: Reallocating budget surplus, balancing resource workloads
6. Strategic Actions & Next Steps: Actionable roadmap for senior leadership to stabilize delivery
"""

    response = client.models.generate_content(
        model="gemini-3.1-flash-lite",
        contents=prompt,
        config={
            "response_mime_type": "application/json"
        }
    )

    # Clean potential markdown wrapping from Gemini output
    text = response.text.strip()
    if text.startswith("```json"):
        text = text[7:]
    if text.endswith("```"):
        text = text[:-3]
    return json.loads(text.strip())


def add_title_slide(prs, title, subtitle):
    # Blank slide layout
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Dark background color for title slide
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRIMARY

    # Title box
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT

    # Subtitle box
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = COLOR_SECONDARY
        p2.alignment = PP_ALIGN.LEFT


def add_bullet_slide(prs, title, subtitle, content):
    # Blank slide layout
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = COLOR_MUTED

    # Bullet content
    bodyBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf_body = bodyBox.text_frame
    tf_body.word_wrap = True

    for i, item in enumerate(content):
        p_item = tf_body.add_paragraph() if i > 0 else tf_body.paragraphs[0]
        p_item.text = f"• {item}"
        p_item.font.size = Pt(16)
        p_item.font.color.rgb = COLOR_DARK
        p_item.space_after = Pt(14)


def add_two_column_slide(prs, title, subtitle, col1_content, col2_content):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = COLOR_MUTED

    # Left Column
    col1Box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.3), Inches(5))
    tf_col1 = col1Box.text_frame
    tf_col1.word_wrap = True
    for i, item in enumerate(col1_content):
        p_item = tf_col1.add_paragraph() if i > 0 else tf_col1.paragraphs[0]
        p_item.text = f"• {item}"
        p_item.font.size = Pt(14)
        p_item.font.color.rgb = COLOR_DARK
        p_item.space_after = Pt(10)

    # Right Column
    col2Box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.3), Inches(5))
    tf_col2 = col2Box.text_frame
    tf_col2.word_wrap = True
    for i, item in enumerate(col2_content):
        p_item = tf_col2.add_paragraph() if i > 0 else tf_col2.paragraphs[0]
        p_item.text = f"• {item}"
        p_item.font.size = Pt(14)
        p_item.font.color.rgb = COLOR_DARK
        p_item.space_after = Pt(10)


def build_presentation(data: dict, output_path: Path):
    prs = Presentation()
    # Use widescreen format (16:9)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    for slide_data in data["slides"]:
        layout = slide_data.get("layout_type", "bullet_slide")
        title = slide_data.get("title", "Executive Slide")
        subtitle = slide_data.get("subtitle", "")

        if layout == "title_slide":
            add_title_slide(prs, title, subtitle)
        elif layout == "two_column_slide":
            col1 = slide_data.get("content", [])
            col2 = slide_data.get("secondary_content", [])
            add_two_column_slide(prs, title, subtitle, col1, col2)
        else:
            content = slide_data.get("content", [])
            add_bullet_slide(prs, title, subtitle, content)

    prs.save(str(output_path))
    print(f"Presentation saved successfully to: {output_path}")


def main():
    output_dir = Path(__file__).parent / "outputs"
    if not output_dir.exists():
        print("Outputs directory does not exist. Run run_all.py first.")
        sys.exit(1)

    reports = []
    for file in output_dir.glob("*_report.json"):
        with open(file, "r") as f:
            reports.append(json.load(f))

    if not reports:
        print("No report files found in outputs/. Run run_all.py first.")
        sys.exit(1)

    print("Synthesizing monthly portfolio metrics via Gemini...")
    try:
        synthesis = get_synthesis(reports)

        ppt_path = output_dir / "Monthly_Executive_Report.pptx"
        build_presentation(synthesis, ppt_path)
    except Exception as e:
        print(f"Error generating presentation: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()
