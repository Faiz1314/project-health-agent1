#!/usr/bin/env python3
"""
Slide Presentation Content Reader

Reads outputs/Monthly_Executive_Report.pptx and prints the slide contents
directly in the terminal in a clean, readable format.
"""

import sys
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("Error: python-pptx package not found. Install it first:")
    print("  pip install python-pptx")
    sys.exit(1)


def read_presentation():
    ppt_path = Path("outputs/Monthly_Executive_Report.pptx")
    if not ppt_path.exists():
        print(f"Error: {ppt_path} not found. Run generate_presentation.py first.")
        return

    prs = Presentation(str(ppt_path))

    print("\n" + "=" * 50)
    print("      GENERATED PPTX SLIDES (TEXT CONTENT)      ")
    print("=" * 50)

    for idx, slide in enumerate(prs.slides):
        print(f"\n📺 SLIDE {idx + 1}")
        print("-" * 50)

        # Separate text findings
        titles = []
        texts = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                txt = shape.text.strip()
                # Try to guess if it's title or body
                if shape.has_text_frame and shape.text_frame.paragraphs[0].font.size:
                    size = shape.text_frame.paragraphs[0].font.size.pt
                    if size >= 24:
                        titles.append(txt)
                        continue
                if txt not in titles:
                    texts.append(txt)

        if titles:
            print(f"Title   : {titles[0]}")
            if len(titles) > 1:
                print(f"Subtitle: {titles[1]}")
        else:
            print("Title   : (No Title)")

        print("\nContent:")
        for t in texts:
            # Format bullets nicely
            lines = t.split("\n")
            for line in lines:
                if line.strip():
                    print(f"  {line}")
        print("-" * 50)

    print("\n" + "=" * 50)


if __name__ == "__main__":
    read_presentation()
